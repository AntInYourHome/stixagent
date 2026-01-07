"""Document loaders for various file formats."""
import os
from typing import List, Dict, Any
from pathlib import Path
import pypdf
from pptx import Presentation
from bs4 import BeautifulSoup
from langchain_core.documents import Document


class DocumentLoader:
    """Load documents from various formats."""
    
    @staticmethod
    def load_pdf(file_path: str) -> List[Document]:
        """Load PDF file and extract text."""
        documents = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        documents.append(Document(
                            page_content=text,
                            metadata={
                                "source": file_path,
                                "page": page_num + 1,
                                "type": "pdf"
                            }
                        ))
        except Exception as e:
            raise ValueError(f"Error loading PDF {file_path}: {str(e)}")
        return documents
    
    @staticmethod
    def load_ppt(file_path: str) -> List[Document]:
        """Load PowerPoint file and extract text."""
        documents = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
                
                if text_parts:
                    text = "\n".join(text_parts)
                    documents.append(Document(
                        page_content=text,
                        metadata={
                            "source": file_path,
                            "slide": slide_num + 1,
                            "type": "ppt"
                        }
                    ))
        except Exception as e:
            raise ValueError(f"Error loading PPT {file_path}: {str(e)}")
        return documents
    
    @staticmethod
    def load_text(file_path: str) -> List[Document]:
        """Load plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                return [Document(
                    page_content=content,
                    metadata={
                        "source": file_path,
                        "type": "text"
                    }
                )]
        except Exception as e:
            raise ValueError(f"Error loading text file {file_path}: {str(e)}")
    
    @staticmethod
    def load_html(file_path: str) -> List[Document]:
        """Load HTML file and extract text."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                text = soup.get_text()
                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = '\n'.join(chunk for chunk in chunks if chunk)
                
                return [Document(
                    page_content=text,
                    metadata={
                        "source": file_path,
                        "type": "html"
                    }
                )]
        except Exception as e:
            raise ValueError(f"Error loading HTML file {file_path}: {str(e)}")
    
    @staticmethod
    def load_document(file_path: str) -> List[Document]:
        """Load document based on file extension."""
        path = Path(file_path)
        extension = path.suffix.lower()
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if extension == '.pdf':
            return DocumentLoader.load_pdf(file_path)
        elif extension in ['.ppt', '.pptx']:
            return DocumentLoader.load_ppt(file_path)
        elif extension in ['.txt', '.md']:
            return DocumentLoader.load_text(file_path)
        elif extension in ['.html', '.htm']:
            return DocumentLoader.load_html(file_path)
        else:
            raise ValueError(f"Unsupported file format: {extension}")

